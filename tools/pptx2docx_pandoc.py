#!/usr/bin/env python3
"""PPTX → Markdown → DOCX (via Pandoc)"""
import sys
sys.path.insert(0, '/Users/yuanweize/Library/Python/3.9/lib/python/site-packages')

import subprocess
import tempfile
from pathlib import Path
from pptx import Presentation


def pptx_to_markdown(input_path):
    """提取 PPTX 文本为 Markdown，按几何位置排序保证阅读顺序"""
    prs = Presentation(input_path)
    md_lines = []
    
    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"# 幻灯片 {i}\n")
        
        # 按 (top, left) 几何排序，保证从上到下、从左到右
        shapes_to_process = []
        for shape in slide.shapes:
            if shape.has_table or (hasattr(shape, "text") and shape.text.strip()):
                top = getattr(shape, 'top', 0) or 0
                left = getattr(shape, 'left', 0) or 0
                shapes_to_process.append((top, left, shape))
        
        shapes_to_process.sort(key=lambda x: (x[0], x[1]))
        
        for _, _, shape in shapes_to_process:
            if hasattr(shape, "text") and shape.text.strip():
                clean_text = shape.text.replace('\x0b', '\n').strip()
                md_lines.append(clean_text + "\n")
            
            if shape.has_table:
                t = shape.table
                try:
                    rows = list(t.rows)
                    if not rows:
                        continue
                    
                    headers = [c.text_frame.text.replace('\n', ' ') for c in rows[0].cells]
                    md_lines.append("| " + " | ".join(headers) + " |")
                    md_lines.append("| " + " | ".join("---" for _ in headers) + " |")
                    
                    for row in rows[1:]:
                        cells = [c.text_frame.text.replace('\n', '<br>') for c in row.cells]
                        md_lines.append("| " + " | ".join(cells) + " |")
                    md_lines.append("")
                except Exception as e:
                    print(f"⚠ 表格转换警告: {e}")
        
        md_lines.append("---\n")
    
    return "\n".join(md_lines)

def ask_overwrite(path):
    """询问是否覆盖已存在的文件"""
    if not path.exists():
        return True
    while True:
        choice = input(f"⚠ {path.name} 已存在，覆盖? [y]是 [N]跳过 [a]全部覆盖: ").lower()
        if choice == '' or choice == 'n':
            return False
        elif choice == 'y':
            return True
        elif choice == 'a':
            return 'all'

def convert(input_path, output_path, md_dir=None, overwrite_all=False):
    md_content = pptx_to_markdown(input_path)
    
    # 检查 docx 是否覆盖
    if not overwrite_all and Path(output_path).exists():
        result = ask_overwrite(Path(output_path))
        if result == False:
            print(f"⏭ 跳过 {output_path}")
            return overwrite_all
        elif result == 'all':
            overwrite_all = True
    
    with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False) as f:
        f.write(md_content)
        md_path = f.name
    
    subprocess.run(['pandoc', md_path, '-o', str(output_path)], check=True)
    
    # 保存 md 文件
    if md_dir:
        md_dir_path = Path(md_dir)
        md_dir_path.mkdir(parents=True, exist_ok=True)  # 自动创建目录
        md_out = md_dir_path / (Path(input_path).stem + '.md')
        if not overwrite_all and md_out.exists():
            result = ask_overwrite(md_out)
            if result == False:
                print(f"⏭ 跳过 {md_out}")
            else:
                if result == 'all':
                    overwrite_all = True
                md_out.write_text(md_content, encoding='utf-8')
                print(f"✓ {md_out}")
        else:
            md_out.write_text(md_content, encoding='utf-8')
            print(f"✓ {md_out}")
    
    Path(md_path).unlink()
    print(f"✓ {output_path}")
    return overwrite_all

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('files', nargs='+')
    parser.add_argument('--md-dir', help='保存 Markdown 的目录')
    parser.add_argument('--out-dir', help='保存 DOCX 的目录')
    parser.add_argument('-y', '--yes', action='store_true', help='自动覆盖，不询问')
    args = parser.parse_args()
    
    overwrite_all = args.yes
    for f in args.files:
        p = Path(f).resolve()
        if args.out_dir:
            out_dir = Path(args.out_dir)
            out_dir.mkdir(parents=True, exist_ok=True)
            out_path = out_dir / (p.stem + '.docx')
        else:
            out_path = p.with_suffix('.docx')
        overwrite_all = convert(p, out_path, args.md_dir, overwrite_all)
